#! /usr/bin/env python

from pptx import Presentation


def main():
	prs = Presentation()
	title_slide_layout = prs.slide_layouts[0]
	slide = prs.slides.add_slide(title_slide_layout)
	title = slide.shapes.title
	subtitle = slide.placeholders[1]

	title.text = "Hello, World!"
	subtitle.text = "python-pptx 可以轻松制作powerpoint!"

	prs.save('../resources/test.pptx')


if __name__ == '__main__':
	main()
